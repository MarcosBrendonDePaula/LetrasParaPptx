import pathlib
import requests
import os
import time
import random
import threading as thread

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt

session = requests.session()

headers = requests.utils.default_headers()
headers.update(
    {
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64; rv:85.0) Gecko/20100101 Firefox/85.0',
    }
)

dlg = None

def loadPage(url=""):
    page = session.get(url, headers=headers)
    return page.text

def getLetra(html):
    try:
        sSoup = BeautifulSoup(html,'html.parser')
        name = sSoup.find(class_="cnt-head_title").find('h1').get_text()
        caixaLetra = sSoup.find(class_='cnt-letra')
        estrofes = caixaLetra.findAll('p')
        letras_br = []
        for l in estrofes:  
            letras_br.append(str(l).replace("<br>","\n").replace("<br/>","\n").replace("</br>","\n").replace("<p>","").replace("</p>","").split('\n')) 
        return [letras_br,name]
    except:
        return ["erro","erro"]

def createPptx(vector,name):
    prs = Presentation()
    left = top = width = height = Inches(2)
    lyt=prs.slide_layouts[0] # choosing a slide layout
    slide = prs.slides.add_slide(lyt) # adding a slide
    title = slide.shapes.title # getTitle
    title.text = name
    lyt = prs.slide_layouts[6] # pegando slide em branco
    for i in vector:
        slide = prs.slides.add_slide(lyt)
        Tbox = slide.shapes.add_textbox(left, top, width, height)
        tframe = Tbox.text_frame
        for x in i:
            P = tframe.add_paragraph()
            P.text = x
            P.font.size = Pt(34)
            pass
    pass
    fname = random.randint(0,8000)
    local = "public/pptx"
    prs.save(str(pathlib.Path(__file__).parent.absolute())+str.format("{0}/{1}.pptx",local,str(fname)+"_"+name))
    return (local,str(fname)+"_"+name)



from flask import Flask,jsonify,render_template,request,redirect,send_file
app = Flask('flaskapp', static_url_path='/', static_folder='public')

def timer(fname):
    minutos = 1
    time.sleep(60*minutos)
    print("Removido:",fname)
    os.remove(str(pathlib.Path(__file__).parent.absolute())+"/public/pptx/{fname}.pptx".format(fname=fname))

@app.route("/gen",methods=['POST'])
def GenPptx():
    dados = getLetra(loadPage(request.form['link']))
    if(dados[0]=="erro"):
        return redirect("/")
    pptx = createPptx(dados[0],dados[1])
    thread.Thread(target=timer,args=(pptx[1],)).start()
    return redirect("pptx/{file}.pptx".format(file=pptx[1]))

@app.route("/",methods=['GET'])
def index():
    return render_template('index.html')

if __name__ == "__main__":
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0',debug=False,port=port)
