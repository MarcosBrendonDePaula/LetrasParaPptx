from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import random

class PresentationGenerator:
    def __init__(self):
        self.default_config = {
            'fontSize': 34,
            'titleFontSize': 44,
            'titleSlideStyle': 'simple',
            'backgroundColor': '#FFFFFF',
            'textColor': '#000000',
            'alignment': 'center',
            'verseSpacing': 1.0,
            'transition': 'none',
            'template': 'blank',
            'fontFamily': 'Calibri',
            'theme': 'default',
            'customTheme': {
                'primary': '#007bff',
                'secondary': '#6c757d',
                'accent': '#28a745'
            },
            'slideMargins': {
                'top': 1,
                'right': 1,
                'bottom': 1,
                'left': 1
            },
            'imageBackground': None,
            'footerText': '',
            'headerText': '',
            'slideNumber': True,
            'verseTransition': 'none',
            'customCSS': ''
        }

    def hex_to_rgb(self, hex_color):
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    def get_alignment(self, align_str):
        alignments = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY
        }
        return alignments.get(align_str.lower(), PP_ALIGN.CENTER)

    def apply_theme(self, slide, config):
        if config['theme'] != 'default':
            # Apply theme-specific styling
            if config['theme'] == 'dark':
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(40, 40, 40)
            elif config['theme'] == 'gradient':
                # Apply gradient background
                background = slide.background.fill
                background.gradient()
                background.gradient_stops[0].color.rgb = self.hex_to_rgb(config['customTheme']['primary'])
                background.gradient_stops[1].color.rgb = self.hex_to_rgb(config['customTheme']['secondary'])
            elif config['theme'] == 'custom':
                # Apply custom theme colors
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = self.hex_to_rgb(config['backgroundColor'])

        # Apply custom image background if provided
        if config['imageBackground']:
            try:
                slide.shapes.add_picture(config['imageBackground'], 0, 0, 
                                       width=slide.slide_width,
                                       height=slide.slide_height)
            except:
                pass

    def create_pptx(self, lyrics_data, config=None):
        if config is None:
            config = {}
        
        # Merge with default config
        config = {**self.default_config, **config}
        
        # Convert colors
        bg_rgb = self.hex_to_rgb(config['backgroundColor'])
        text_rgb = self.hex_to_rgb(config['textColor'])
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        self.apply_theme(title_slide, config)
        
        # Title
        title = title_slide.shapes.title
        title.text = lyrics_data['title']
        title_font = title.text_frame.paragraphs[0].font
        title_font.size = Pt(float(config['titleFontSize']))
        title_font.color.rgb = RGBColor(*text_rgb)
        title_font.name = config['fontFamily']

        # Subtitle (Artist)
        if config['titleSlideStyle'] in ['withArtist', 'full']:
            subtitle = title_slide.placeholders[1]
            subtitle.text = lyrics_data['artist']
            subtitle_font = subtitle.text_frame.paragraphs[0].font
            subtitle_font.color.rgb = RGBColor(*text_rgb)
            subtitle_font.name = config['fontFamily']

        # Content slides
        blank_layout = prs.slide_layouts[6]
        alignment = self.get_alignment(config['alignment'])
        
        # Calculate margins
        margins = config['slideMargins']
        left = Inches(float(margins['left']))
        top = Inches(float(margins['top']))
        width = prs.slide_width - (Inches(float(margins['left']) + float(margins['right'])))
        height = prs.slide_height - (Inches(float(margins['top']) + float(margins['bottom'])))

        for verse_lines in lyrics_data['verses']:
            slide = prs.slides.add_slide(blank_layout)
            self.apply_theme(slide, config)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            
            # Add header if specified
            if config['headerText']:
                header = text_frame.add_paragraph()
                header.text = config['headerText']
                header.font.size = Pt(12)
                header.font.color.rgb = RGBColor(*text_rgb)
                header.alignment = PP_ALIGN.CENTER
            
            # Add lines with proper formatting
            for i, line in enumerate(verse_lines):
                paragraph = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
                paragraph.text = line
                paragraph.font.size = Pt(float(config['fontSize']))
                paragraph.font.color.rgb = RGBColor(*text_rgb)
                paragraph.font.name = config['fontFamily']
                paragraph.alignment = alignment
                paragraph.line_spacing = Pt(float(config['fontSize']) * float(config['verseSpacing']))

            # Add footer if specified
            if config['footerText']:
                footer = text_frame.add_paragraph()
                footer.text = config['footerText']
                footer.font.size = Pt(12)
                footer.font.color.rgb = RGBColor(*text_rgb)
                footer.alignment = PP_ALIGN.CENTER

            # Add slide number if enabled
            if config['slideNumber']:
                slide_num = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.7),
                                                   Inches(1), Inches(0.3))
                slide_num.text_frame.paragraphs[0].text = str(slide.slide_id)
                slide_num.text_frame.paragraphs[0].font.size = Pt(10)
                slide_num.text_frame.paragraphs[0].font.color.rgb = RGBColor(*text_rgb)

        # Generate unique filename
        fname = random.randint(0, 8000)
        file_name = f"{fname}_{lyrics_data['title'].replace(' ', '_')}"
        file_path = f"public/pptx/{file_name}.pptx"
        
        # Save presentation
        prs.save(file_path)
        return file_name
